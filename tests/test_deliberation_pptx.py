"""심의 현황팩 PPTX 산출 테스트 (CLAUDE.md §8.13 C4·C5).

QuotaAssessment → A3 pptx 구조 검증. 위성 타일 fetch 는 monkeypatch 로 스킵(표만).
"""

import io

from pptx import Presentation
from pptx.util import Cm

from app.services import deliberation_pptx as dp
from app.schemas.quota import QuotaAssessment, QuotaResult, FacilityQuota
from app.schemas.survey import SurveyResult, SurveyDong, FacilityCategory, SurveyFacility


def _assessment():
    return QuotaAssessment(
        address="서울 동작구 본동 441", site_sgg="11590", site_lat=37.5, site_lon=127.0,
        radius=1000, ym="202606", gu_infant=7994, gu_households=188064,
        survey=SurveyResult(
            address="a", site_sgg="11590", site_lat=37.5, site_lon=127.0, radius=1000, ym="202606",
            dongs=[SurveyDong(name="노량진1동", ratio=0.9896, total_pop=33349, total_hh=18457,
                              applied_pop=33001, applied_hh=18264, same_sgg=True),
                   SurveyDong(name="이촌2동", ratio=0.2222, total_pop=7733, total_hh=3570,
                              applied_pop=1718, applied_hh=793, same_sgg=False, flagged=True)],
            applied_pop_total=33001, applied_hh_total=18264),
        facilities=[
            FacilityCategory(category="작은도서관", count=2, capacity=None, items=[
                SurveyFacility(name="가도서관", addr="가로1", dist_m=272, lat=37.501, lon=127.001),
                SurveyFacility(name="나도서관", addr="나로2", dist_m=338, lat=37.499, lon=127.0)]),
            FacilityCategory(category="경로당", count=1, items=[
                SurveyFacility(name="다경로당", addr="", dist_m=79, lat=37.5005, lon=127.0, src="vworld")]),
            FacilityCategory(category="어린이집", count=1, capacity=2974, capacity_scope="동작구", items=[
                SurveyFacility(name="라어린이집", addr="라로4", dist_m=102, lat=37.4995, lon=127.0)]),
        ],
        results=[QuotaResult(label="", new_households=981, applied_households=18264, facilities=[
            FacilityQuota(name="작은도서관", households=19245, expected_people=1154.7,
                          required_area=1886.6, existing_area=1000.16, planned_area=515.92,
                          legal_min=158, legal_min_confidence="high", verdict="부족시설",
                          plan_ok=True, plan_diff=357.92),
            FacilityQuota(name="경로당", households=18264, required_area=-1189.0,
                          existing_area=3015.63, legal_min=330, legal_min_confidence="high",
                          verdict="충족시설")])])


def test_build_pptx_structure(monkeypatch):
    monkeypatch.setattr(dp, "_basemap", lambda lat, lon, radius, client: None)  # 타일 fetch 스킵
    data = dp.build_pptx(_assessment())
    assert data[:2] == b"PK"                       # zip(pptx) 시그니처

    prs = Presentation(io.BytesIO(data))
    assert abs(prs.slide_width - Cm(42.0)) < 1000 and abs(prs.slide_height - Cm(29.7)) < 1000
    slides = list(prs.slides)
    # 걸침표 1 + 시설(도서관·경로당) 2 + 총량제 1 = 4
    assert len(slides) == 4
    # 각 슬라이드에 표 최소 1개
    for s in slides:
        assert any(sh.has_table for sh in s.shapes)

    # 걸침표: 계 행에 적용세대
    t0 = next(sh.table for sh in slides[0].shapes if sh.has_table)
    joined = " ".join(c.text for r in t0.rows for c in r.cells)
    assert "18,264" in joined and "노량진1동" in joined and "⚠" in joined  # 타시군구 플래그

    # 총량제 슬라이드: 부족/충족 판정 텍스트 (표 셀 안)
    tq = []
    for s in slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                tq.append(sh.text_frame.text)
            if sh.has_table:
                tq += [c.text for r in sh.table.rows for c in r.cells]
    allt = " ".join(tq)
    assert "부족시설" in allt and "충족시설" in allt


def test_build_pptx_map_when_basemap_ok(monkeypatch):
    """위성 타일이 되면 도서관 슬라이드에 그림(배경) + 도형(핀) 이 붙는다."""
    from PIL import Image

    def fake_basemap(lat, lon, radius, client):
        return Image.new("RGB", (dp._MAP_PX, dp._MAP_PX), (30, 30, 30)), 15, dp._MAP_PX / 2, dp._MAP_PX / 2
    monkeypatch.setattr(dp, "_basemap", fake_basemap)
    monkeypatch.setattr(dp.tiles, "latlon_to_global_px", lambda la, lo, z: (dp._MAP_PX / 2, dp._MAP_PX / 2))
    monkeypatch.setattr(dp.tiles, "meters_to_pixels", lambda m, lat, z: 100.0)

    data = dp.build_pptx(_assessment())
    prs = Presentation(io.BytesIO(data))
    lib_slide = list(prs.slides)[1]  # 도서관
    assert any("PICTURE" in str(sh.shape_type) for sh in lib_slide.shapes)  # 위성 배경
    assert any(sh.name.startswith("반경") for sh in lib_slide.shapes)        # 편집가능 반경원

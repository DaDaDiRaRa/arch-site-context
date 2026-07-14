"""C7 주변현황도 PPTX + 라우터 테스트 (CLAUDE.md §8.13)."""

import io

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Cm

from app.main import app
import app.services.surroundings_pptx as spptx
import app.routers.surroundings as srouter
from app.schemas.surroundings import SurroundingsResult, SurroundCategory, SurroundItem

client = TestClient(app)


def _result():
    return SurroundingsResult(
        address="서울 동작구 본동 441", site_lat=37.5, site_lon=127.0, radius=1000,
        ring_radii=[250, 500, 750],
        categories=[
            SurroundCategory(name="교통", count=3, color=(233, 30, 99), items=[
                SurroundItem(name="노들역", dist_m=120, lat=37.501, lon=127.0)]),
            SurroundCategory(name="교육", count=6, color=(33, 150, 243), items=[
                SurroundItem(name="영본초등학교", dist_m=300, lat=37.499, lon=127.0)]),
        ],
        narrative="대상지 반경 1000m 내 지하철 노들역 인접 · 교육시설 6개소 등.")


def test_build_surroundings_pptx(monkeypatch):
    # 위성 타일 실패 → 표·서술문만 (network 스킵)
    def boom(*a, **k):
        raise RuntimeError("no tiles")
    monkeypatch.setattr(spptx.tiles, "compose_basemap", boom)

    data = spptx.build_surroundings_pptx(_result())
    assert data[:2] == b"PK"
    prs = Presentation(io.BytesIO(data))
    assert abs(prs.slide_width - Cm(42.0)) < 1000
    slides = list(prs.slides)
    assert len(slides) == 1
    s = slides[0]
    tbl = next(sh.table for sh in s.shapes if sh.has_table)
    joined = " ".join(c.text for r in tbl.rows for c in r.cells)
    assert "교통" in joined and "교육" in joined and "노들역" in joined
    # 서술문 텍스트 존재
    txt = " ".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
    assert "주변현황" in txt and "노들역" in txt


def test_surroundings_endpoints(monkeypatch):
    monkeypatch.setattr(srouter, "collect_surroundings", lambda addr, radius: _result())
    monkeypatch.setattr(spptx, "build_surroundings_pptx", lambda r, client=None: b"PK\x03\x04x")

    r1 = client.post("/surroundings", json={"address": "서울 동작구 본동 441", "radius": 1000})
    assert r1.status_code == 200
    assert r1.json()["categories"][0]["name"] == "교통"
    assert "노들역" in r1.json()["narrative"]

    r2 = client.post("/surroundings/pptx", json={"address": "서울 동작구 본동 441", "radius": 1000})
    assert r2.status_code == 200
    assert r2.json()["url"].endswith(".pptx")
    assert r2.json()["categories"]["교육"] == 6
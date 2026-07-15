"""종합 대지 읽기 PPT — 아키타입·해석·의견·드라이버·교차·POR·방법론 (S4 기본 포함).

/board 결과(BoardResult dict)를 **대지분석 덱과 같은 디자인 언어**(deck.style: 다크프레임·
대괄호타이틀·PANEL·레드캡션밴드)로 A3 편집가능 PPTX 로 조립. 새 숫자 0 — 값은 board 그대로.
① 사실 종합(해석)과 ② 종합 의견은 **벽으로 분리**(green/amber edge)해 사실과 의견을 구분한다.
"""
from __future__ import annotations

import io
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Pt

import app.deck.style as k

HL = k.RGBColor(0xFF, 0xE0, 0x66)
GREEN = k.RGBColor(0x4C, 0xC0, 0x5B)   # ① 사실 벽
AMBER = k.RGBColor(0xE8, 0x8A, 0x1E)   # ② 의견 벽


def _txt(v, n=1600):
    s = (v or "").strip()
    return s[:n]


def _n(v):
    try:
        f = float(v)
        return f"{int(f):,}" if f == int(f) else f"{f:,.1f}"
    except (TypeError, ValueError):
        return str(v or "-")


UP = k.RGBColor(0xE8, 0x8A, 0x1E)     # 상회
DOWN = k.RGBColor(0x35, 0x9F, 0xE0)   # 하회
PINK = k.RGBColor(0xF0, 0x2E, 0x8A)   # 결론 핵심어 강조


def _runs_tb(sl, x, y, w, h, runs, *, size, align=None, bold=True, anchor=None):
    """한 문단 안에 여러 색 run — 헤드라인 문장(핵심어 강조)용."""
    b = sl.shapes.add_textbox(x, y, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align if align is not None else k.PP_ALIGN.LEFT
    try:
        p.line_spacing = 1.15
    except Exception:  # noqa: BLE001
        pass
    for text, color in runs:
        r = p.add_run(); r.text = text
        r.font.size, r.font.name, r.font.bold, r.font.color.rgb = Pt(size), k.F, bold, color
    return b


# ── 종합 결론 (헤드라인 문장 + 핵심어 강조 — 발표용, 데이터에서 조립·LLM 0) ──
def slide_conclusion(prs, board):
    arch = board.get("archetype") or {}
    drivers = board.get("design_drivers") or []
    cross = board.get("cross_implications") or []
    facts = [f for f in (board.get("facts") or []) if isinstance(f.get("index"), (int, float))]

    findings = []  # (statement_runs, evidence)
    if arch.get("name"):
        ev = " · ".join(f"{e.get('key')} {e.get('detail')}" for e in (arch.get("evidence") or [])[:2])
        findings.append(([("이 동네는 ", k.WHITE), (arch["name"], PINK), (" 유형이다.", k.WHITE)], ev))
    if drivers:
        d = drivers[0]
        ev = " · ".join([d.get("response", "")] + [f"{e.get('key')} {e.get('detail')}" for e in (d.get("evidence") or [])[:1]])
        findings.append(([("설계의 지배 변수는 ", k.WHITE), (d.get("name", ""), PINK), ("다.", k.WHITE)], ev))
    if facts:
        f = max(facts, key=lambda x: abs((x.get("index") or 100) - 100))
        band = f.get("index_band") or "비슷"
        findings.append(([(f.get("item", ""), PINK),
                          (f" {_n(f.get('value'))}{f.get('unit', '')} — 전국 대비 {band}.", k.WHITE)],
                         f"지수 {f.get('index')} (전국=100)"))
    if cross:
        c = cross[0]
        findings.append(([(c.get("name", ""), PINK), (" 검토가 필요하다.", k.WHITE)], _txt(c.get("text"), 80)))
    if not findings:
        return False

    sl = k.blank_slide(prs)
    # 라벨 칩 (중앙 상단)
    k.rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, Cm(18.0), Cm(1.3), Cm(6.0), Cm(1.0), fill=PINK, alpha_pct=90)
    k.tb(sl, Cm(18.0), Cm(1.3), Cm(6.0), Cm(1.0), "종합 결론", size=13, color=k.WHITE, bold=True,
         align=k.PP_ALIGN.CENTER, anchor=k.MSO_ANCHOR.MIDDLE)
    k.tb(sl, Cm(1.5), Cm(2.7), Cm(39), Cm(1.3), "이 필지는 무엇이 다른가?", size=28, color=k.WHITE,
         bold=True, align=k.PP_ALIGN.CENTER)

    y = 5.3
    step = min(4.6, (23.0 - y) / max(1, len(findings[:4])))
    for runs, ev in findings[:4]:
        _runs_tb(sl, Cm(3.5), Cm(y), Cm(35), Cm(1.6), runs, size=21, align=k.PP_ALIGN.CENTER)
        k.tb(sl, Cm(3.5), Cm(y + 1.55), Cm(35), Cm(0.9), ev, size=11, color=k.MUTE, align=k.PP_ALIGN.CENTER)
        y += step

    k.caption_band(sl, [("데이터에서 뽑은 핵심 결론 ", k.WHITE, True),
                        ("· 근거는 뒤 슬라이드(지표·드라이버·교차)에서", HL, True)])
    return True


# ── 동네 프로필 (지표 전국=100 발산형 막대차트) ──
def slide_indicators(prs, board):
    if not [f for f in (board.get("facts") or []) if isinstance(f.get("index"), (int, float))]:
        return False
    sl = k.blank_slide(prs)
    k.bracket_title(sl, "동네 프로필", "PROFILE · 주요 지표 전국=100 대비 (상회 ▶ / 하회 ◀)")
    k.index_chart(sl, board.get("facts"), top=4.9)
    k.caption_band(sl, [("주요 지표를 전국 평균(100) 대비 막대로 ", k.WHITE, True),
                        ("· 상회/하회는 방향일 뿐 판정 아님(참고)", HL, True)])
    return True


def _sentences(text: str) -> list:
    """문단 → 문장 리스트. 마침표+공백 뒤에서 분리 (괄호 안 소수점 '45.1%' 는 공백 없어 안전)."""
    s = (text or "").strip()
    if not s:
        return []
    parts = re.split(r"(?<=[.!?。])\s+", s)
    return [p.strip() for p in parts if p.strip()]


def _prose(sl, x, y, w, h, text, *, accent, size=15):
    """긴 서술을 **문장 단위 불릿 + 여백·행간**으로 — 한 덩어리 벽 대신 읽기 좋게."""
    sents = _sentences(text) or [_txt(text)]
    b = sl.shapes.add_textbox(x, y, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    for i, s in enumerate(sents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(11)
        try:
            p.line_spacing = 1.2
        except Exception:  # noqa: BLE001 — 일부 환경 미지원 시 기본 행간
            pass
        r1 = p.add_run(); r1.text = "▪  "
        r1.font.size, r1.font.name, r1.font.bold, r1.font.color.rgb = Pt(size), k.F, True, accent
        r2 = p.add_run(); r2.text = s
        r2.font.size, r2.font.name, r2.font.color.rgb = Pt(size), k.F, k.WHITE
    return b


# ── 커버 (아키타입 헤드라인) ──
def _cover(prs, address, board):
    arch = board.get("archetype") or {}
    sl = k.blank_slide(prs)
    k.rect(sl, MSO_SHAPE.RECTANGLE, 0, Cm(11.0), k.A3_W, Cm(0.16), fill=k.RED)
    k.tb(sl, Cm(1.5), Cm(2.8), Cm(39), Cm(2), "종합 대지 읽기", size=40, color=k.WHITE, bold=True)
    k.tb(sl, Cm(1.5), Cm(5.6), Cm(39), Cm(1), "SITE SYNTHESIS", size=15, color=k.MUTE)
    name = arch.get("name")
    if name:
        k.tb(sl, Cm(1.5), Cm(7.3), Cm(39), Cm(1.4), f"이 동네는 — {name}", size=27, color=HL, bold=True)
        sub = f"{arch.get('group', '')} · {_txt(arch.get('description'), 70)}"
        k.tb(sl, Cm(1.5), Cm(9.0), Cm(39), Cm(0.9), sub, size=13, color=k.WHITE)
    k.tb(sl, Cm(1.5), Cm(12.2), Cm(39), Cm(1), address, size=15, color=k.WHITE)
    k.tb(sl, Cm(1.5), Cm(13.0), Cm(39), Cm(1),
         "아키타입 · 종합해석 · 종합의견 · 설계드라이버 · 교차시사점 · 프로그램 함의(POR)",
         size=12, color=k.MUTE)


# ── 종합 개요 (커버리지 + 아키타입 근거) ──
def slide_overview(prs, board):
    cov = board.get("coverage") or []
    arch = board.get("archetype") or {}
    sl = k.blank_slide(prs)
    k.bracket_title(sl, "종합 개요", "OVERVIEW · 도메인 커버리지 + 동네 유형 근거")
    # 커버리지 카드 (있음/확인불가)
    cw, ch, gx = Cm(7.6), Cm(3.2), Cm(0.5)
    for i, c in enumerate(cov[:5]):
        x = Cm(1.3) + i * (cw + gx)
        ok = c.get("available")
        edge = GREEN if ok else k.MUTE
        k.rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(3.6), cw, ch, fill=k.PANEL)
        k.rect(sl, MSO_SHAPE.RECTANGLE, x, Cm(3.85), Cm(0.12), ch - Cm(0.5), fill=edge)
        k.tb(sl, x + Cm(0.4), Cm(3.9), cw - Cm(0.6), Cm(0.6), c.get("domain", ""), size=12, color=k.WHITE, bold=True)
        k.tb(sl, x + Cm(0.4), Cm(4.6), cw - Cm(0.6), Cm(0.6), "확보" if ok else "확인 불가",
             size=11, color=edge, bold=True)
        k.tb(sl, x + Cm(0.4), Cm(5.2), cw - Cm(0.6), Cm(1.4), _txt(c.get("detail"), 40), size=9, color=k.MUTE)
    # 아키타입 근거
    if arch:
        k.panel_head(sl, Cm(1.3), Cm(7.6), f"동네 유형 · {arch.get('name', '')}", w=Cm(30), size=15)
        y = 9.1
        for e in (arch.get("evidence") or [])[:6]:
            k.rect(sl, MSO_SHAPE.OVAL, Cm(1.3), Cm(y + 0.1), Cm(0.24), Cm(0.24), fill=k.RED)
            prox = f" ({e.get('proximity')})" if e.get("proximity") else ""
            k.tb(sl, Cm(1.9), Cm(y), Cm(38), Cm(0.7), f"{e.get('key')} — {e.get('detail')}{prox}",
                 size=12, color=k.WHITE)
            y += 0.9
        alts = arch.get("alternatives") or []
        if alts:
            k.tb(sl, Cm(1.9), Cm(y + 0.2), Cm(38), Cm(0.7), "혼재 특성(차점): " + ", ".join(alts[:3]),
                 size=11, color=k.MUTE)
    k.caption_band(sl, [("데이터 확보 현황과 동네 유형 근거 ", k.WHITE, True),
                        ("· 확인불가는 숨기지 않음(참고)", HL, True)])


# ── 종합 해석·의견 (한 슬라이드 좌우 분리: ① 사실 / ② 의견) ──
def _syn_panel(sl, x, y, w, h, chip, edge, body):
    k.rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=k.PANEL)
    k.rect(sl, MSO_SHAPE.RECTANGLE, x, y, Cm(0.2), h, fill=edge)
    k.rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x + Cm(0.6), y + Cm(0.5), Cm(7.6), Cm(0.95), fill=edge, alpha_pct=90)
    k.tb(sl, x + Cm(0.6), y + Cm(0.5), Cm(7.6), Cm(0.95), chip, size=12, color=k.WHITE, bold=True,
         align=k.PP_ALIGN.CENTER, anchor=k.MSO_ANCHOR.MIDDLE)
    _prose(sl, x + Cm(0.8), y + Cm(1.9), w - Cm(1.5), h - Cm(2.3), body, accent=edge, size=15)


def slide_synthesis(prs, board):
    syn = board.get("synthesis") or {}
    itp = syn.get("interpretation")
    jdg = syn.get("judgment")
    if not itp and not jdg:
        return False
    sl = k.blank_slide(prs)
    k.bracket_title(sl, "종합 해석 · 의견", "SYNTHESIS · ① 검증된 사실 / ② 종합 의견(참고·최종 판단은 사람)")
    # 내용(문장 수)에 맞춰 패널 높이 산정 후 세로 중앙 — 큰 빈 박스 방지
    ns = max([len(_sentences(b)) for b in (itp, jdg) if b] + [1])
    hc = min(20.0, max(6.5, 3.0 + ns * 2.6))
    top = 3.6 + max(0.0, (23.7 - hc) / 2)
    W, H = Cm(19.3), Cm(hc)
    if itp:
        _syn_panel(sl, Cm(1.3), Cm(top), W, H, "① 검증된 사실", GREEN, itp)
    if jdg:
        _syn_panel(sl, Cm(21.4), Cm(top), W, H, "② 종합 의견 · 참고", AMBER, jdg)
    k.caption_band(sl, [("① 검증된 사실 · ② 그 위의 종합 의견 ", k.WHITE, True),
                        ("· 최종 판단은 사람", HL, True)])
    return True


# ── 설계 드라이버 ──
def _center_top(rows, row_h, top=3.6, bottom=27.3):
    """카드 그리드를 세로 중앙에 — 위에 몰리고 아래 비는 것 방지."""
    return top + max(0.0, (bottom - top - rows * row_h) / 2)


def slide_drivers(prs, board):
    drivers = board.get("design_drivers") or []
    if not drivers:
        return False
    sl = k.blank_slide(prs)
    k.bracket_title(sl, "설계 드라이버", "DESIGN DRIVERS · 지배 드라이버 (증거강도 랭킹 · 검토 신호)")
    maxs = max((d.get("strength") or 0 for d in drivers), default=1) or 1
    y = _center_top(len(drivers[:5]), 3.75)
    for d in drivers[:5]:
        k.rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.3), Cm(y), Cm(39.4), Cm(3.5), fill=k.PANEL)
        k.tb(sl, Cm(1.7), Cm(y + 0.35), Cm(2.6), Cm(1), f"#{d.get('rank')}", size=22, color=k.RED, bold=True)
        k.tb(sl, Cm(4.2), Cm(y + 0.4), Cm(24), Cm(1), d.get("name", ""), size=16, color=k.WHITE, bold=True)
        # 강도 바
        bx, bw = Cm(30.0), 8.0
        k.rect(sl, MSO_SHAPE.RECTANGLE, bx, Cm(y + 0.55), Cm(bw), Cm(0.45), fill=k.NAVY)
        k.rect(sl, MSO_SHAPE.RECTANGLE, bx, Cm(y + 0.55), Cm(bw * (d.get("strength") or 0) / maxs), Cm(0.45), fill=k.RED)
        k.tb(sl, bx, Cm(y + 1.05), Cm(bw), Cm(0.5), f"증거강도 {float(d.get('strength') or 0):.1f}", size=9, color=k.MUTE)
        k.tb(sl, Cm(4.2), Cm(y + 1.45), Cm(25), Cm(0.9), d.get("response", ""), size=12, color=k.MUTE)
        ev = " · ".join(f"{e.get('key')} {e.get('detail')}" for e in (d.get("evidence") or [])[:2])
        k.tb(sl, Cm(4.2), Cm(y + 2.4), Cm(35), Cm(0.9), f"근거: {ev}", size=9.5, color=k.MUTE)
        y += 3.75
    k.caption_band(sl, [("분석 → 설계 다리 ", k.WHITE, True),
                        ("· 검토 신호이며 판정·제안안 아님(참고)", HL, True)])
    return True


# 도메인·카테고리 색 팔레트
DOMAIN_COL = {"인구": DOWN, "수급": k.RGBColor(0x2A, 0xA1, 0x98), "재해": UP}
CARD_PALETTE = [
    k.RGBColor(0x35, 0x9F, 0xE0), k.RGBColor(0x2A, 0xA1, 0x98), k.RGBColor(0xE8, 0x8A, 0x1E),
    k.RGBColor(0x4C, 0xC0, 0x5B), k.RGBColor(0x9B, 0x6F, 0xD4), k.RGBColor(0xC8, 0x81, 0x3C),
]


def _chip(sl, x, y, text, col):
    w = Cm(0.9 + 0.42 * len(text))
    k.rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Cm(0.72), fill=col, alpha_pct=92)
    k.tb(sl, x, y, w, Cm(0.72), text, size=9.5, color=k.WHITE, bold=True,
         align=k.PP_ALIGN.CENTER, anchor=k.MSO_ANCHOR.MIDDLE)
    return w


# ── 교차시사점 (도메인 교차 카드 그리드) ──
def slide_cross(prs, board):
    cross = board.get("cross_implications") or []
    if not cross:
        return False
    sl = k.blank_slide(prs)
    k.bracket_title(sl, "교차시사점", "CROSS-CONTEXT · 도메인 교차 (인구 × 수급 × 재해)")
    # 항목 수에 맞춰 카드 높이를 키워 세로를 채움 (2열, 적을수록 큼)
    W = Cm(19.3)
    xs = [Cm(1.3), Cm(21.4)]
    cc = cross[:6]
    rows = (len(cc) + 1) // 2
    ch = min(8.2, max(4.2, (23.4 - (rows - 1) * 0.5) / rows))
    RH = ch + 0.5
    y0 = 3.6  # 상단 정렬 — 위부터 채워 자연스럽게
    for i, c in enumerate(cc):
        x, y = xs[i % 2], Cm(y0 + (i // 2) * RH)
        k.rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, W, Cm(ch), fill=k.PANEL)
        # 도메인 칩 (인구 × 수급 …)
        cxp = x + Cm(0.5)
        for j, dm in enumerate((c.get("domains") or [])[:3]):
            if j > 0:
                k.tb(sl, cxp, y + Cm(0.5), Cm(0.55), Cm(0.7), "×", size=12, color=k.MUTE, bold=True,
                     align=k.PP_ALIGN.CENTER, anchor=k.MSO_ANCHOR.MIDDLE)
                cxp += Cm(0.55)
            cxp += _chip(sl, cxp, y + Cm(0.5), dm, DOMAIN_COL.get(dm, k.MUTE)) + Cm(0.2)
        k.tb(sl, x + Cm(0.5), y + Cm(1.6), W - Cm(0.9), Cm(0.8), c.get("name", ""), size=15, color=k.WHITE, bold=True)
        k.tb(sl, x + Cm(0.5), y + Cm(2.55), W - Cm(0.9), Cm(1.4), _txt(c.get("text"), 90), size=12, color=k.WHITE)
        basis = " · ".join(f"{b.get('key')} {b.get('detail')}" for b in (c.get("basis") or [])[:2])
        k.tb(sl, x + Cm(0.5), y + Cm(ch - 0.85), W - Cm(0.9), Cm(0.6), f"근거: {basis}", size=9.5, color=k.MUTE)
    k.caption_band(sl, [("도메인을 가로지르는 참고 시사점 ", k.WHITE, True),
                        ("· 규칙 조합(LLM 0)·판단 아님", HL, True)])
    return True


# ── 프로그램 함의(POR) — 카테고리 컬러 카드 그리드 ──
def slide_program(prs, board):
    por = board.get("program_implications") or []
    if not por:
        return False
    groups = {}
    for p in por:
        groups.setdefault(p.get("category", "기타"), []).append(p)
    sl = k.blank_slide(prs)
    k.bracket_title(sl, "프로그램 함의 (POR)", "PROGRAM · 카테고리별 공간·프로그램 권고 (검토 체크리스트)")
    W, CH = Cm(12.9), 6.6
    xs = [Cm(1.3), Cm(14.55), Cm(27.8)]
    cats = list(groups.items())[:6]
    y_top = _center_top((len(cats) + 2) // 3, CH)
    for i, (cat, items) in enumerate(cats):
        col = CARD_PALETTE[i % len(CARD_PALETTE)]
        x = xs[i % 3]
        yb = y_top + (i // 3) * CH  # float cm (Cm() 로 감쌀 기준값)
        k.rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(yb), W, Cm(CH - 0.4), fill=k.PANEL)
        k.rect(sl, MSO_SHAPE.RECTANGLE, x, Cm(yb + 0.3), Cm(0.14), Cm(CH - 1.0), fill=col)
        k.tb(sl, x + Cm(0.45), Cm(yb + 0.28), W - Cm(0.7), Cm(0.8), cat, size=13, color=col, bold=True)
        yy = yb + 1.3
        for p in items[:4]:
            k.rect(sl, MSO_SHAPE.OVAL, x + Cm(0.5), Cm(yy + 0.14), Cm(0.2), Cm(0.2), fill=col)
            k.tb(sl, x + Cm(0.9), Cm(yy), W - Cm(1.25), Cm(1.15), _txt(p.get("recommendation"), 46),
                 size=10.5, color=k.WHITE)
            yy += 1.2
    k.caption_band(sl, [("맥락 → 건축 프로그램 번역 ", k.WHITE, True),
                        ("· 검토 권고이며 최종 결정은 사람(참고)", HL, True)])
    return True


# ── 컨셉·설계방향 제안 (발표 커버 — AI 제안·격리, services/concept.py 참조) ──
def slide_concept(prs, board):
    """이미지 레퍼런스(PROJECT VALUE)처럼 컨셉명 + 키워드 3개로 설계 방향 제안.

    board['concept'] 가 있을 때만 렌더(opt-in). 이건 S4 원칙의 의도적 예외(창작)라
    '① 검증된 사실'이 아니라 'AI 제안'임을 라벨·핑크 강조로 사실 슬라이드와 분리한다.
    """
    con = board.get("concept") or {}
    kws = con.get("keywords") or []
    if not kws:
        return False
    sl = k.blank_slide(prs)
    # 상단 라벨 칩 — 이건 '사실'이 아니라 'AI 제안'
    k.rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, Cm(16.5), Cm(1.4), Cm(9.0), Cm(1.0), fill=AMBER, alpha_pct=92)
    k.tb(sl, Cm(16.5), Cm(1.4), Cm(9.0), Cm(1.0), "AI 제안 · 설계 방향", size=13, color=k.WHITE, bold=True,
         align=k.PP_ALIGN.CENTER, anchor=k.MSO_ANCHOR.MIDDLE)
    k.tb(sl, Cm(1.5), Cm(2.9), Cm(39), Cm(1.0), "PROJECT DIRECTION", size=15, color=k.MUTE,
         bold=True, align=k.PP_ALIGN.CENTER)

    # 컨셉명 (창작·핑크). 규칙 폴백이면 name 이 비어 '설계 방향' 을 대신 씀.
    name = (con.get("name") or "").strip()
    headline = f"– {name} –" if name else "설계 방향"
    _runs_tb(sl, Cm(1.5), Cm(4.6), Cm(39), Cm(3.2), [(headline, PINK)], size=54,
             align=k.PP_ALIGN.CENTER, anchor=k.MSO_ANCHOR.MIDDLE)
    tag = (con.get("tagline") or "").strip()
    if tag:
        k.tb(sl, Cm(1.5), Cm(7.9), Cm(39), Cm(0.9), tag, size=14, color=k.WHITE, align=k.PP_ALIGN.CENTER)

    # 키워드 행 (word · word · word) — 이미지의 대형 중앙 배열
    words = [str(w.get("word", "")).strip() for w in kws if str(w.get("word", "")).strip()]
    runs = []
    for i, w in enumerate(words):
        if i > 0:
            runs.append((" · ", k.MUTE))
        runs.append((w, k.WHITE))
    _runs_tb(sl, Cm(1.5), Cm(9.6), Cm(39), Cm(2.0), runs, size=34, align=k.PP_ALIGN.CENTER,
             anchor=k.MSO_ANCHOR.MIDDLE)

    # 키워드 풀이 (뜻 + 근거) — 각 키워드가 어느 드라이버/사실에서 나왔는지 추적
    y = 12.4
    for w in kws:
        word = str(w.get("word", "")).strip()
        gloss = _txt(w.get("gloss"), 60)
        _runs_tb(sl, Cm(4.0), Cm(y), Cm(34), Cm(0.9),
                 [(f"{word}: ", PINK), (gloss, k.WHITE)], size=15, align=k.PP_ALIGN.CENTER)
        basis = _txt(w.get("basis"), 70)
        if basis:
            k.tb(sl, Cm(4.0), Cm(y + 0.85), Cm(34), Cm(0.7), f"근거: {basis}", size=10,
                 color=k.MUTE, align=k.PP_ALIGN.CENTER)
        y += 1.85

    label = con.get("label") or ""
    k.caption_band(sl, [("AI 설계방향 제안 · 컨셉·네이밍은 참고안 ", k.WHITE, True),
                        ("· 검증/재현 보장 없음 · 최종 컨셉·제안서는 사람", HL, True)])
    return True


# ── 방법론·출처 부록 ──
def slide_methodology(prs, board):
    meth = board.get("methodology") or {}
    sources = meth.get("sources") or []
    if not sources:
        return False
    sl = k.blank_slide(prs)
    k.bracket_title(sl, "방법론·출처", "METHODOLOGY · 사용 출처 (공모·감사 대비)")
    rows = []
    for s in sources[:14]:
        rows.append([s.get("name", ""), s.get("id", "") or s.get("source", ""),
                     _txt(s.get("note") or s.get("detail"), 40)])
    if rows:
        k.table(sl, Cm(1.3), Cm(3.6), Cm(39.4), ["출처", "식별자", "비고"], rows,
                ratios=[3.4, 3.0, 3.6], fs=9.5)
    k.caption_band(sl, _txt(meth.get("summary"), 100) or "이 보드 수치의 출처·산정식·한계.")
    return True


# ── 조립 ──
def build_board_deck(board: dict) -> bytes:
    """BoardResult dict → 종합 대지 읽기 A3 PPTX bytes. 데이터 없는 섹션은 graceful skip."""
    address = (board.get("site") or {}).get("address") or ""
    prs = Presentation()
    prs.slide_width, prs.slide_height = k.A3_W, k.A3_H
    _cover(prs, address, board)
    slide_conclusion(prs, board)
    slide_overview(prs, board)
    slide_indicators(prs, board)
    slide_synthesis(prs, board)
    slide_drivers(prs, board)
    slide_cross(prs, board)
    slide_program(prs, board)
    slide_concept(prs, board)  # opt-in — board['concept'] 있을 때만 (AI 제안·격리)
    slide_methodology(prs, board)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

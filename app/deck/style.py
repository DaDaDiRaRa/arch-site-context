"""대지분석 보드 공유 디자인·유틸 모듈.

지도 슬라이드 생성기(광역도·조망분석·입지현황·용도)와 데이터 팩트 슬라이드가 공유.
디자인 원칙: 다크 네이비 프레임 · 컬러 위성(흑백 안 함) · 대괄호 타이틀 ·
하단 레드 캡션밴드(핵심어 강조) · 빨간 SITE · N 나침반 · 네이티브 편집가능.
"""
from __future__ import annotations

import io
from typing import Optional

from PIL import Image, ImageEnhance
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt
from lxml import etree
from pyproj import Transformer

from app.services.tiles import latlon_to_global_px as _latlon_to_global_px

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
A3_W, A3_H = Cm(42.0), Cm(29.7)

# ── 색 토큰 ──
NAVY = RGBColor(0x1B, 0x24, 0x38); PANEL = RGBColor(0x26, 0x33, 0x4D)
RED = RGBColor(0xE1, 0x24, 0x2B); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1E, 0x1E, 0x1E); MUTE = RGBColor(0xB4, 0xBF, 0xD2)
F = "맑은 고딕"
# 지하철 노선색
LINE_COLOR = {"1": RGBColor(0x00, 0x52, 0xA4), "2": RGBColor(0x00, 0xA8, 0x4D),
              "3": RGBColor(0xEF, 0x7C, 0x1C), "4": RGBColor(0x00, 0xA4, 0xE3),
              "5": RGBColor(0x8B, 0x3D, 0xA8), "6": RGBColor(0xB6, 0x50, 0x0A),
              "7": RGBColor(0x74, 0x7F, 0x00), "9": RGBColor(0xBD, 0xB0, 0x00)}
# 건물 높이색
def hcol(h: float) -> RGBColor:
    if h < 15: return RGBColor(0x4C, 0xC0, 0x5B)
    if h < 35: return RGBColor(0xF5, 0xC5, 0x18)
    if h < 60: return RGBColor(0xF3, 0x8B, 0x1E)
    return RGBColor(0xE8, 0x3A, 0x2F)

# 조망 유형 (건물 최고높이 기준)
VIEW = {
    "HIGH": (RGBColor(0xF0, 0x3B, 0x2E), "HIGH & BUSY CITY VIEW", "고층 시가지 · 차폐"),
    "MID":  (RGBColor(0xF5, 0xA6, 0x23), "MIDRISE CITY VIEW", "중층 시가지"),
    "LOW":  (RGBColor(0x35, 0x9F, 0xE0), "LOWRISE VIEW", "저층 · 개방"),
    "OPEN": (RGBColor(0x4C, 0xC0, 0x5B), "OPEN / GREEN VIEW", "개방 · 조망"),
}
def view_of(h: float) -> str:
    return "OPEN" if h < 8 else ("LOW" if h < 25 else ("MID" if h < 60 else "HIGH"))

# ── EPSG:5186 (arch-site-model 좌표계) ──
_TO_4326 = Transformer.from_crs("EPSG:5186", "EPSG:4326", always_xy=True)
_TO_5186 = Transformer.from_crs("EPSG:4326", "EPSG:5186", always_xy=True)

def local_to_latlon(lx: float, ly: float, ox: float, oy: float) -> tuple[float, float]:
    """model 로컬미터(+origin_offset) → (lat, lon)."""
    lon, lat = _TO_4326.transform(lx + ox, ly + oy)
    return lat, lon

def latlon_to_local(lat: float, lon: float, ox: float, oy: float) -> tuple[float, float]:
    """(lat,lon) → model 로컬미터(=5186-origin). 사이트 로컬좌표 산출용."""
    x, y = _TO_5186.transform(lon, lat)
    return x - ox, y - oy


# ── 텍스트·도형 헬퍼 ──
def tb(sl, x, y, w, h, s, *, size=12, color=INK, bold=False, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, rot=0, wrap=True):
    b = sl.shapes.add_textbox(x, y, w, h); b.rotation = rot
    tf = b.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    if not wrap:  # 짧은 라벨(핀 번호 등) 강제 1줄 — 좁은 박스 자동 줄바꿈 방지
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(s if isinstance(s, list) else [s]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = align
        r = p.add_run(); r.text = str(ln)
        r.font.size, r.font.name, r.font.bold, r.font.color.rgb = Pt(size), F, bold, color
    return b


def alpha(shape, pct):
    sf = shape.fill._xPr.find(f'.//{A}solidFill'); clr = sf.find(f'{A}srgbClr')
    etree.SubElement(clr, f'{A}alpha').set('val', str(int(pct * 1000)))


def rect(sl, shp, x, y, w, h, *, fill=None, line=None, lw=Pt(1), alpha_pct=None, dash=None):
    s = sl.shapes.add_shape(shp, x, y, w, h)
    if fill is None: s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
        if alpha_pct is not None: alpha(s, alpha_pct)
    if line is None: s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = lw
        if dash:
            ln = s.line._get_or_add_ln(); etree.SubElement(ln, f'{A}prstDash').set('val', dash)
    s.shadow.inherit = False
    return s


# ── 컴포넌트 ──
def prep_satellite(png: bytes) -> bytes:
    """컬러 유지 + 대비/명도 살짝 (흑백 안 함 — 사용자 지시)."""
    img = Image.open(io.BytesIO(png)).convert("RGB")
    img = ImageEnhance.Contrast(img).enhance(1.03)
    o = io.BytesIO(); img.save(o, "PNG"); return o.getvalue()


def dark_frame(sl):
    rect(sl, MSO_SHAPE.RECTANGLE, 0, 0, A3_W, A3_H, fill=NAVY)


def bracket_title(sl, title, subtitle="", x=Cm(1.3), y=Cm(0.5)):
    tb(sl, x, y, Cm(30), Cm(1.1), f"[ {title} ]", size=22, color=WHITE, bold=True)
    if subtitle:
        tb(sl, x, y + Cm(1.25), Cm(30), Cm(0.7), subtitle, size=11, color=MUTE)


def add_map(sl, png, x, y, w, size):
    """위성 배치 + canvas px→EMU 변환기 반환. (ex, ey, scale)"""
    sl.shapes.add_picture(io.BytesIO(png), x, y, w, w)
    scale = int(w) / size
    def ex(c): return Emu(int(x) + int(c * scale))
    def ey(c): return Emu(int(y) + int(c * scale))
    return ex, ey, scale


def to_canvas(lat, lon, z, mcx, mcy, size):
    """(lat,lon) → basemap canvas px (add_map ex/ey에 넣기 전)."""
    gx, gy = _latlon_to_global_px(lat, lon, z)
    return gx - mcx + size / 2, gy - mcy + size / 2


def site_pill(sl, ex, ey, cx, cy, label="SITE"):
    rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, ex(cx) - Cm(1.1), ey(cy) - Cm(0.55), Cm(2.2), Cm(1.1),
         fill=RED, line=WHITE, lw=Pt(2))
    tb(sl, ex(cx) - Cm(1.1), ey(cy) - Cm(0.55), Cm(2.2), Cm(1.1), label, size=14, color=WHITE,
       bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def parcel_canvas(geom, z, mcx, mcy, size):
    """law parcel_geometry(MultiPolygon/Polygon, WGS84) → 최대 외곽링의 canvas px 목록."""
    if not geom:
        return []
    t, coords = geom.get("type"), geom.get("coordinates") or []
    rings = []
    if t == "MultiPolygon":
        rings = [poly[0] for poly in coords if poly]
    elif t == "Polygon":
        rings = [coords[0]] if coords else []
    if not rings:
        return []
    ring = max(rings, key=len)
    from app.services.tiles import latlon_to_global_px as _latlon_to_global_px
    out = []
    for pt in ring:
        gx, gy = _latlon_to_global_px(pt[1], pt[0], z)  # [lon,lat]
        out.append((gx - mcx + size / 2, gy - mcy + size / 2))
    return out


def site_marker(sl, ex, ey, cx, cy, parcel_pts=None, label="SITE"):
    """SITE 표기 — 필지 형상(있으면) + 빨간 십자, 없으면 빨간 알약."""
    if parcel_pts and len(parcel_pts) >= 3:
        try:
            fb = sl.shapes.build_freeform(int(ex(parcel_pts[0][0])), int(ey(parcel_pts[0][1])), scale=1.0)
            fb.add_line_segments([(int(ex(x)), int(ey(y))) for x, y in parcel_pts[1:]], close=True)
            shp = fb.convert_to_shape()
            shp.fill.solid(); shp.fill.fore_color.rgb = RED; alpha(shp, 26)
            shp.line.color.rgb = WHITE; shp.line.width = Pt(2.25); shp.shadow.inherit = False
        except Exception:
            return site_pill(sl, ex, ey, cx, cy, label)
        ccx = sum(p[0] for p in parcel_pts) / len(parcel_pts)
        ccy = sum(p[1] for p in parcel_pts) / len(parcel_pts)
        cr = Cm(0.55)
        cx_e, cy_e = int(ex(ccx)), int(ey(ccy))
        c = sl.shapes.add_shape(MSO_SHAPE.CROSS, cx_e - cr // 2, cy_e - cr // 2, cr, cr)
        c.fill.solid(); c.fill.fore_color.rgb = RED; c.line.color.rgb = WHITE; c.line.width = Pt(1); c.shadow.inherit = False
        tb(sl, Emu(cx_e) - Cm(1.5), Emu(cy_e) + Cm(0.35), Cm(3), Cm(0.6), label, size=11, color=WHITE,
           bold=True, align=PP_ALIGN.CENTER)
    else:
        site_pill(sl, ex, ey, cx, cy, label)


def n_compass(sl, x, y):
    tb(sl, x, y, Cm(1.5), Cm(1.0), "N", size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def caption_band(sl, runs):
    """하단 레드 캡션밴드. runs = str 또는 [(text, color|None, bold)] (핵심어 컬러강조)."""
    rect(sl, MSO_SHAPE.RECTANGLE, 0, A3_H - Cm(1.9), A3_W, Cm(1.9), fill=RED)
    b = sl.shapes.add_textbox(Cm(1.3), A3_H - Cm(1.75), Cm(39.4), Cm(1.5))
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if isinstance(runs, str):
        runs = [(runs, None, True)]
    for text, color, bold in runs:
        r = p.add_run(); r.text = text
        r.font.size, r.font.name, r.font.bold = Pt(13), F, bold
        r.font.color.rgb = color if color else WHITE
    return b


def dot_label(sl, ex, ey, cx, cy, lines, *, dot=WHITE, dot_line=INK, box=INK, box_alpha=72,
              size=8.5, dot_r=Cm(0.13)):
    """건물/역 라벨: 흰 점 + 어두운 박스 텍스트(다중행)."""
    rect(sl, MSO_SHAPE.OVAL, Emu(int(ex(cx))) - dot_r, Emu(int(ey(cy))) - dot_r, dot_r * 2, dot_r * 2,
         fill=dot, line=dot_line, lw=Pt(0.75))
    txt = lines if isinstance(lines, list) else [lines]
    longest = max((len(t) for t in txt), default=4)
    w = Cm(max(2.6, 0.34 * longest + 1.0))
    hh = Cm(0.52 * len(txt) + 0.1)
    rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, Emu(int(ex(cx))) + Cm(0.22), Emu(int(ey(cy))) - hh / 2,
         w, hh, fill=box, alpha_pct=box_alpha)
    tb(sl, Emu(int(ex(cx))) + Cm(0.22), Emu(int(ey(cy))) - hh / 2, w, hh, txt, size=size,
       color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── 데이터 표현 컴포넌트 (팩트 슬라이드 공용 — 지도와 같은 디자인 언어) ──
ROW_ALT = RGBColor(0x20, 0x2B, 0x42)  # 표 줄무늬(PANEL·NAVY 사이 톤)


def blank_slide(prs):
    """다크 프레임 빈 슬라이드 (레이아웃 6=빈)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    dark_frame(sl)
    return sl


def panel_head(sl, x, y, title, *, w=Cm(14), size=16, color=WHITE):
    """레드 탭 + 흰 제목 (슬라이드 내 소제목)."""
    rect(sl, MSO_SHAPE.RECTANGLE, x, y, Cm(0.14), Cm(1.0), fill=RED)
    tb(sl, x + Cm(0.4), y - Cm(0.05), w, Cm(1.0), title, size=size, color=color, bold=True)


def kpi_card(sl, x, y, w, h, label, value, sub="", *, accent=RED):
    """PANEL 카드 + 큰 숫자 (대지개요·생활맥락 지표용)."""
    rect(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=PANEL)
    rect(sl, MSO_SHAPE.RECTANGLE, x, y + Cm(0.25), Cm(0.12), h - Cm(0.5), fill=accent)
    tb(sl, x + Cm(0.45), y + Cm(0.3), w - Cm(0.7), Cm(0.6), label, size=10.5, color=MUTE, bold=True)
    tb(sl, x + Cm(0.45), y + Cm(0.92), w - Cm(0.7), Cm(1.15), str(value), size=21, color=WHITE, bold=True)
    if sub:
        tb(sl, x + Cm(0.45), y + h - Cm(0.72), w - Cm(0.7), Cm(0.6), sub, size=9, color=MUTE)


def _cell(tbl, r, c, text, *, color=WHITE, bold=False, size=9.0, fill=NAVY, align=PP_ALIGN.LEFT):
    cell = tbl.cell(r, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.margin_left = Cm(0.18); cell.margin_right = Cm(0.1)
    cell.margin_top = Cm(0.02); cell.margin_bottom = Cm(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size, run.font.name, run.font.bold, run.font.color.rgb = Pt(size), F, bold, color


def table(sl, x, y, w, headers, rows, *, ratios=None, fs=9.0, header_fs=9.5,
          row_h=Cm(0.82), header_h=Cm(0.92)):
    """네이티브 편집가능 표 (다크 테마). headers=list[str], rows=list[list].

    첫 열 좌측정렬·나머지 가운데. 셀 채움을 직접 지정해 기본 표스타일 위에 다크 톤을 덮는다.
    """
    ncols = len(headers)
    nrows = len(rows) + 1
    total_h = header_h + row_h * len(rows)
    gfx = sl.shapes.add_table(nrows, ncols, x, y, w, total_h)
    tbl = gfx.table
    tbl.first_row = False
    tbl.horz_banding = False
    if ratios:
        tot = float(sum(ratios))
        for i, r in enumerate(ratios):
            tbl.columns[i].width = Emu(int(int(w) * r / tot))
    tbl.rows[0].height = header_h
    for r in range(1, nrows):
        tbl.rows[r].height = row_h
    for c, h in enumerate(headers):
        _cell(tbl, 0, c, h, color=WHITE, bold=True, size=header_fs, fill=PANEL, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows, start=1):
        fill = NAVY if ri % 2 else ROW_ALT
        for c, val in enumerate(row):
            al = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            _cell(tbl, ri, c, "" if val is None else val, fill=fill, align=al, size=fs)
    return gfx

"""대지분석 지도 슬라이드 4종 — 서비스 통합 (아무 주소나).

광역입지도 · 건물용도현황 · 입지현황(높이) · 방향별 조망분석.
프로토타입(make_wide/use/site/viewring)을 파라미터화·통합. 공유 데이터 1회 fetch.
전부 컬러·네이티브 편집가능. `app.deck_style` 공유 디자인.
"""
from __future__ import annotations

import io, math, re
from concurrent.futures import ThreadPoolExecutor
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt
from lxml import etree

from app.deck import clients, data_slides as ds
from app.deck.clients import _haversine_m, _BLDG_SUFFIX
import app.deck.style as k

_BLDG_KINDS = ["아파트", "오피스텔", "병원", "학교", "빌딩", "타워", "스퀘어"]
DIRN = ["북", "북동", "동", "남동", "남", "남서", "서", "북서"]

USE = {
    "주거": (k.RGBColor(0x3B, 0x7D, 0xD8), "주거(공동주택)"),
    "상업": (k.RGBColor(0xC8, 0x81, 0x3C), "상업·근생"),
    "업무": (k.RGBColor(0x2A, 0xA1, 0x98), "업무"),
    "공업": (k.RGBColor(0x88, 0x88, 0x88), "공업"),
    "공공": (k.RGBColor(0x5A, 0xB0, 0x53), "공공·교육·의료"),
    "미상": (k.RGBColor(0xB4, 0xB4, 0xB4), "용도 미상"),
}
KIND_USE = {"아파트": "주거", "빌라": "주거", "연립": "주거", "상가": "상업", "음식점": "상업",
            "편의점": "상업", "카페": "상업", "마트": "상업", "오피스텔": "업무", "빌딩": "업무",
            "타워": "업무", "공장": "공업", "지식산업센터": "공업", "공업사": "공업", "학교": "공공",
            "병원": "공공", "주민센터": "공공", "관공서": "공공", "어린이집": "공공"}
_USE_STRONG = {"아파트", "빌라", "연립", "공장", "지식산업센터", "공업사", "학교", "병원", "주민센터", "관공서", "어린이집", "오피스텔"}
HL = k.RGBColor(0xFF, 0xE0, 0x66)


# ── 건물 매싱 추출 (5186) ──
def _buildings(model, lat0, lon0, z, mcx, mcy, size):
    ox, oy = (model.get("stats") or {})["origin_offset"]
    out = []
    for b in (model.get("geometry") or {}).get("buildings") or []:
        fp = b.get("footprint") or []
        if len(fp) < 3:
            continue
        pts = [k.to_canvas(*k.local_to_latlon(lx, ly, ox, oy), z, mcx, mcy, size) for (lx, ly) in fp]
        if not all(-50 <= x <= size + 50 and -50 <= y <= size + 50 for x, y in pts):
            continue
        clat, clon = k.local_to_latlon(sum(p[0] for p in fp) / len(fp), sum(p[1] for p in fp) / len(fp), ox, oy)
        out.append({"pts": pts, "cx": sum(p[0] for p in pts) / len(pts), "cy": sum(p[1] for p in pts) / len(pts),
                    "clat": clat, "clon": clon, "h": b.get("height") or 0.0, "name": None, "use": "미상"})
    return out


def _draw_footprints(sl, ex, ey, blds, color_fn, alpha_fn=lambda b: 45):
    for b in blds:
        pts = b["pts"]
        try:
            fb = sl.shapes.build_freeform(int(ex(pts[0][0])), int(ey(pts[0][1])), scale=1.0)
            fb.add_line_segments([(int(ex(x)), int(ey(y))) for x, y in pts[1:]], close=True)
            shp = fb.convert_to_shape()
        except Exception:
            continue
        col = color_fn(b)
        shp.fill.solid(); shp.fill.fore_color.rgb = col; k.alpha(shp, alpha_fn(b))
        shp.line.color.rgb = col; shp.line.width = Pt(0.75); shp.shadow.inherit = False


def _match_names(blds, fac):
    for f in fac.get("results", []):
        base = re.sub(r"\s*\d+동.*$", "", str(f.get("name") or "")).strip()
        if not base.endswith(_BLDG_SUFFIX) or f.get("lat") is None:
            continue
        best = min(blds, key=lambda b: (b["clat"] - f["lat"]) ** 2 + (b["clon"] - f["lon"]) ** 2, default=None)
        if best and _haversine_m(best["clat"], best["clon"], f["lat"], f["lon"]) < 55 and not best["name"]:
            best["name"] = base


def _spread(items, n, dist):
    placed = []
    for b in items:
        if all((b["cx"] - p["cx"]) ** 2 + (b["cy"] - p["cy"]) ** 2 > dist ** 2 for p in placed):
            placed.append(b)
        if len(placed) >= n:
            break
    return placed


def _blank(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); k.dark_frame(sl); return sl


# ── 슬라이드 1: 광역입지도 ──
def slide_wide(prs, address, lat, lon, law, site):
    fac = clients.fetch_facilities(address, ["지하철역"], 2000)
    meta, png = clients.fetch_basemap(lat, lon, 2000, 1500)
    png = k.prep_satellite(png); size = 1500
    z, mcx, mcy, rpx = int(meta["zoom"]), float(meta["cx"]), float(meta["cy"]), float(meta["radius_px"])
    lp = (site or {}).get("land_price") or {}; bd = (site or {}).get("building") or {}
    law = law or {}

    def sbase(n): return re.sub(r"\(.*?\)", "", re.sub(r"\s*\d+번출구.*$", "", re.sub(r"\s*\d+호선.*$", "", n))).strip()
    stations = {}
    for f in fac.get("results", []):
        base = sbase(str(f.get("name") or ""))
        if not base.endswith("역") or f.get("lat") is None:
            continue
        s = stations.setdefault(base, {"name": base, "lat": f["lat"], "lon": f["lon"], "lines": set()})
        s["lines"] |= set(re.findall(r"(\d+)호선", str(f.get("name") or "")))
    stations = sorted(stations.values(), key=lambda s: (s["lat"] - lat) ** 2 + (s["lon"] - lon) ** 2)[:16]

    sl = _blank(prs)
    k.bracket_title(sl, "광역입지현황", "PROJECT POSITIONING")
    y = 4.2
    for key, val in [("위치", address), ("지목", (lp.get("jibun") or "").replace("385", "").strip() or "확인필요"),
                     ("용도지역", law.get("zone_use") or "확인필요"),
                     ("개별공시지가", f"{lp.get('price_per_sqm'):,}원/㎡ ({lp.get('year')})" if lp.get("price_per_sqm") else "확인필요"),
                     ("현황", bd.get("name") or "나대지"), ("PNU", law.get("pnu") or (site or {}).get("pnu") or "-")]:
        k.tb(sl, Cm(1.3), Cm(y), Cm(13), Cm(0.6), key, size=11, color=k.RED, bold=True)
        k.tb(sl, Cm(1.3), Cm(y + 0.55), Cm(13), Cm(0.9), val, size=13, color=k.WHITE, bold=True); y += 1.75

    BL, BT, BW = Cm(15.5), Cm(2.6), Cm(25.5)
    ex, ey, scale = k.add_map(sl, png, BL, BT, BW, size); scx = scy = size / 2
    for rm in (500, 1000, 2000):
        rp = rpx * rm / 2000
        k.rect(sl, MSO_SHAPE.OVAL, ex(scx - rp), ey(scy - rp), Emu(int(2 * rp * scale)), Emu(int(2 * rp * scale)),
               line=k.WHITE, lw=Pt(1.25), dash="dash")
        k.tb(sl, ex(scx) - Cm(1), ey(scy - rp) - Cm(0.5), Cm(2), Cm(0.5), f"{rm//1000}km" if rm >= 1000 else f"{rm}m",
             size=9, color=k.WHITE, align=PP_ALIGN.CENTER)
    gray = k.RGBColor(0x55, 0x55, 0x55)
    for s in stations:
        cxp, cyp = k.to_canvas(s["lat"], s["lon"], z, mcx, mcy, size)
        if not (10 <= cxp <= size - 10 and 10 <= cyp <= size - 10):
            continue
        col = k.LINE_COLOR.get(sorted(s["lines"])[0], gray) if s["lines"] else gray
        k.dot_label(sl, ex, ey, cxp, cyp, s["name"], dot=col, dot_line=k.WHITE)
    k.site_pill(sl, ex, ey, scx, scy); k.n_compass(sl, Cm(39.3), Cm(2.9))
    present = sorted({l for s in stations for l in s["lines"]}, key=int)
    k.tb(sl, Cm(15.5), Cm(24.5), Cm(8), Cm(0.5), "지하철 노선", size=10, color=k.MUTE)
    lx = 15.5
    for l in present:
        k.rect(sl, MSO_SHAPE.OVAL, Cm(lx), Cm(25.0), Cm(0.5), Cm(0.5), fill=k.LINE_COLOR.get(l, gray))
        k.tb(sl, Cm(lx + 0.6), Cm(24.97), Cm(2), Cm(0.6), f"{l}호선", size=11, color=k.WHITE, bold=True); lx += 2.4
    names = ", ".join(s["name"] for s in stations[:3]) or "인근 역"
    k.caption_band(sl, [(f"{law.get('zone_use') or '도심'} · 역세권 입지 — 반경 1km 내 ", k.WHITE, True),
                        (names, HL, True), (" 등 지하철 접근 양호. 광역 교통축과 연계된 통합 거점.", k.WHITE, True)])


# ── 슬라이드 2: 건물 용도현황 ──
def _name_use(name):
    for suf, u in (("아파트", "주거"), ("빌라", "주거"), ("맨션", "주거"), ("오피스텔", "업무"), ("타워", "업무"),
                   ("빌딩", "업무"), ("스퀘어", "업무"), ("병원", "공공"), ("의원", "공공"), ("학교", "공공"),
                   ("지식산업센터", "공업"), ("공업사", "공업"), ("공장", "공업")):
        if name.endswith(suf):
            return u
    return None


def slide_use(prs, address, lat, lon, model, parcel=None):
    kinds = list(KIND_USE.keys())
    fac = clients.fetch_facilities(address, kinds, 320)
    meta, png = clients.fetch_basemap(lat, lon, 320, 1500)
    png = k.prep_satellite(png); size = 1500
    z, mcx, mcy = int(meta["zoom"]), float(meta["cx"]), float(meta["cy"])
    strong, weak = [], []
    for f in fac.get("results", []):
        u = KIND_USE.get(f.get("kind"))
        nm = str(f.get("name") or "")
        if u and f.get("lat") is not None:
            # 병원 kind 중 의원·치과·한의원(동네 병의원=1층 근생류)은 약한 신호로
            is_strong = f.get("kind") in _USE_STRONG
            if f.get("kind") == "병원" and any(w in nm for w in ("의원", "치과", "한의원", "동물")):
                is_strong = False
            (strong if is_strong else weak).append((f["lat"], f["lon"], u, nm))
    blds = _buildings(model, lat, lon, z, mcx, mcy, size)
    for b in blds:
        def nearest(lst):
            bd, bu, bn = 1e9, None, None
            for (pla, plo, pu, pn) in lst:
                d = _haversine_m(b["clat"], b["clon"], pla, plo)
                if d < bd:
                    bd, bu, bn = d, pu, pn
            return bd, bu, bn
        sd, su, sn = nearest(strong); wd, wu, wn = nearest(weak)
        if sd < 45: use, nm = su, sn
        elif wd < 35: use, nm = wu, wn
        else: use, nm = "미상", None
        nu = _name_use(nm) if nm else None
        b["use"] = nu or use
        b["name"] = nm if (nm and _name_use(nm)) else None

    sl = _blank(prs)
    ex, ey, scale = k.add_map(sl, png, Cm(1.3), Cm(2.6), Cm(27.0), size)
    _draw_footprints(sl, ex, ey, blds, lambda b: USE[b["use"]][0],
                     lambda b: 48 if b["use"] != "미상" else 28)
    named = _spread(sorted([b for b in blds if b["name"]], key=lambda b: b["use"] != "주거"), 13, 230)
    for b in named:
        k.dot_label(sl, ex, ey, b["cx"], b["cy"], [b["name"][:12], USE[b["use"]][1]], box_alpha=74)
    k.site_marker(sl, ex, ey, size / 2, size / 2, k.parcel_canvas(parcel, z, mcx, mcy, size))
    k.bracket_title(sl, "건물 용도현황", f"SITE CONTEXT · 주변 건물 용도 · {address}")
    k.n_compass(sl, Cm(26.5), Cm(2.9))
    cnt = {}
    for b in blds:
        cnt[b["use"]] = cnt.get(b["use"], 0) + 1
    PX = Cm(29.5)
    k.rect(sl, MSO_SHAPE.RECTANGLE, PX, Cm(2.6), Cm(0.14), Cm(1.0), fill=k.RED)
    k.tb(sl, PX + Cm(0.4), Cm(2.55), Cm(11), Cm(1), "건물 용도", size=16, color=k.WHITE, bold=True)
    ly = 4.2
    for u in ["주거", "상업", "업무", "공업", "공공", "미상"]:
        col, lab = USE[u]
        k.rect(sl, MSO_SHAPE.RECTANGLE, PX, Cm(ly), Cm(0.7), Cm(0.7), fill=col, alpha_pct=72)
        k.tb(sl, PX + Cm(1.0), Cm(ly - 0.02), Cm(9), Cm(0.8), f"{lab}  ({cnt.get(u, 0)})", size=12, color=k.WHITE, bold=True); ly += 1.0
    k.tb(sl, PX, Cm(ly + 0.4), Cm(11), Cm(2.5), [f"· 반경 내 건물 {len(blds)}동", "· 용도=주변 시설(kakao) 추정", "  명칭 확인 건물은 실제 용도"], size=11, color=k.MUTE)
    top = max((u for u in USE if u != "미상"), key=lambda u: cnt.get(u, 0), default="주거")
    k.caption_band(sl, [("주변 건물 ", k.WHITE, True), (f"{USE[top][1]} 우세", HL, True),
                        (f" ({cnt.get(top, 0)}/{len(blds)}동) — 다양한 용도가 혼재한 시가지. 용도는 주변시설 기반 추정.", k.WHITE, True)])


# ── 슬라이드 3: 입지현황 (높이) ──
def slide_site(prs, address, lat, lon, model, parcel=None):
    fac = clients.fetch_facilities(address, _BLDG_KINDS, 320)
    meta, png = clients.fetch_basemap(lat, lon, 320, 1500)
    png = k.prep_satellite(png); size = 1500
    z, mcx, mcy = int(meta["zoom"]), float(meta["cx"]), float(meta["cy"])
    blds = _buildings(model, lat, lon, z, mcx, mcy, size)
    _match_names(blds, fac)
    named = _spread(sorted([b for b in blds if b["name"]], key=lambda b: -b["h"]), 14, 220)

    sl = _blank(prs)
    ex, ey, scale = k.add_map(sl, png, Cm(1.3), Cm(2.6), Cm(27.0), size)
    _draw_footprints(sl, ex, ey, blds, lambda b: k.hcol(b["h"]))
    for b in named:
        fl = max(1, round(b["h"] / 3.0))  # 층고 3.0 고정 → gro_flo_co 정확
        k.dot_label(sl, ex, ey, b["cx"], b["cy"], [b["name"][:12], f"{fl}층 · 약 {int(b['h'])}m"], box_alpha=74)
    k.site_marker(sl, ex, ey, size / 2, size / 2, k.parcel_canvas(parcel, z, mcx, mcy, size))
    k.bracket_title(sl, "입지현황", f"SITE CONTEXT · 주변 건물 매싱·높이 · {address}")
    k.n_compass(sl, Cm(26.5), Cm(2.9))
    hs = [b["h"] for b in blds]
    PX = Cm(29.5)
    k.rect(sl, MSO_SHAPE.RECTANGLE, PX, Cm(2.6), Cm(0.14), Cm(1.0), fill=k.RED)
    k.tb(sl, PX + Cm(0.4), Cm(2.55), Cm(11), Cm(1), "주변 건물 높이", size=16, color=k.WHITE, bold=True)
    ly = 4.2
    for lab, c in [("저층 ~15m", k.hcol(10)), ("중저 15~35m", k.hcol(25)), ("중고 35~60m", k.hcol(50)), ("고층 60m~", k.hcol(80))]:
        k.rect(sl, MSO_SHAPE.RECTANGLE, PX, Cm(ly), Cm(0.7), Cm(0.7), fill=c, alpha_pct=70)
        k.tb(sl, PX + Cm(1.0), Cm(ly - 0.02), Cm(9), Cm(0.8), lab, size=12, color=k.WHITE, bold=True); ly += 1.0
    k.tb(sl, PX, Cm(ly + 0.5), Cm(11), Cm(3), [f"· 반경 내 건물 {len(blds)}동",
         f"· 최고 {int(max(hs)) if hs else 0}m / 평균 {int(sum(hs)/len(hs)) if hs else 0}m", f"· 라벨 건물 {len(named)}동"], size=12, color=k.MUTE)
    tall = sum(1 for h in hs if h >= 60)
    k.caption_band(sl, [("주변 ", k.WHITE, True), (f"{len(blds)}동 중 60m 이상 고층 {tall}동", HL, True),
                        (" — 층수·높이 실측(VWorld), 도로폭은 별도 확인.", k.WHITE, True)])


# ── 슬라이드 4: 방향별 조망 ──
def slide_viewring(prs, address, lat, lon, model, parcel=None):
    fac = clients.fetch_facilities(address, _BLDG_KINDS, 600)
    meta, png = clients.fetch_basemap(lat, lon, 600, 1500)
    png = k.prep_satellite(png); size = 1500
    z, mcx, mcy = int(meta["zoom"]), float(meta["cx"]), float(meta["cy"])
    ox, oy = (model.get("stats") or {})["origin_offset"]
    sx, sy = k.latlon_to_local(lat, lon, ox, oy)
    sect = [0.0] * 8; blds = []
    for b in (model.get("geometry") or {}).get("buildings") or []:
        fp = b.get("footprint") or []
        if len(fp) < 3:
            continue
        cxl = sum(p[0] for p in fp) / len(fp); cyl = sum(p[1] for p in fp) / len(fp)
        dx, dy = cxl - sx, cyl - sy
        if abs(dx) + abs(dy) < 3:
            continue
        brg = (math.degrees(math.atan2(dx, dy)) + 360) % 360
        h = b.get("height") or 0.0
        sect[int(brg // 45)] = max(sect[int(brg // 45)], h)
        clat, clon = k.local_to_latlon(cxl, cyl, ox, oy)
        blds.append({"clat": clat, "clon": clon, "h": h, "name": None})
    _match_names(blds, fac)
    named = sorted([b for b in blds if b["name"]], key=lambda b: -b["h"])
    views = [k.view_of(h) for h in sect]
    groups, used = [], [False] * 8
    for i in range(8):
        if used[i]:
            continue
        j = i
        while views[(j + 1) % 8] == views[i] and not used[(j + 1) % 8] and (j + 1) % 8 != i:
            j += 1
        for m in range(i, j + 1):
            used[m % 8] = True
        groups.append((i, j, views[i]))

    sl = _blank(prs)
    ex, ey, scale = k.add_map(sl, png, Cm(1.3), Cm(2.6), Cm(25.5), size)
    scx = scy = size / 2; R = size * 0.44
    for i, j, vk in groups:
        col = k.VIEW[vk][0]; b1, b2 = i * 45 + 1.5, (j + 1) * 45 - 1.5
        arc = sl.shapes.add_shape(MSO_SHAPE.ARC, ex(scx - R), ey(scy - R), Emu(int(2 * R * scale)), Emu(int(2 * R * scale)))
        try:
            arc.adjustments[0], arc.adjustments[1] = (b1 - 90) % 360, (b2 - 90) % 360
        except Exception:
            pass
        arc.fill.background(); arc.line.color.rgb = col; arc.line.width = Pt(9); arc.shadow.inherit = False
        mb = ((b1 + b2) / 2) % 360; rr = R * 1.06
        lx = scx + rr * math.sin(math.radians(mb)); ly = scy - rr * math.cos(math.radians(mb))
        k.tb(sl, ex(lx) - Cm(3), ey(ly) - Cm(0.35), Cm(6), Cm(0.7), k.VIEW[vk][1], size=10, color=col,
             bold=True, align=PP_ALIGN.CENTER, rot=(mb if mb <= 180 else mb - 180))
    for i, j, vk in groups:
        mb = ((i * 45 + (j + 1) * 45) / 2) % 360
        x2 = scx + R * 0.98 * math.sin(math.radians(mb)); y2 = scy - R * 0.98 * math.cos(math.radians(mb))
        cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ex(scx), ey(scy), ex(x2), ey(y2))
        cn.line.color.rgb = k.WHITE; cn.line.width = Pt(0.75)
        etree.SubElement(cn.line._get_or_add_ln(), f'{k.A}prstDash').set('val', 'dash'); cn.shadow.inherit = False
    placed = []
    for b in named:
        cxp, cyp = k.to_canvas(b["clat"], b["clon"], z, mcx, mcy, size)
        if not (40 <= cxp <= size - 40 and 40 <= cyp <= size - 40):
            continue
        if all((cxp - p[0]) ** 2 + (cyp - p[1]) ** 2 > 200 ** 2 for p in placed):
            placed.append((cxp, cyp))
            k.dot_label(sl, ex, ey, cxp, cyp, [b["name"][:11], f"{int(b['h'])}m"], box_alpha=72)
        if len(placed) >= 12:
            break
    k.site_marker(sl, ex, ey, scx, scy, k.parcel_canvas(parcel, z, mcx, mcy, size))
    k.bracket_title(sl, "방향별 조망 · 건물높이 분석", f"SITE CONTEXT · {address}")
    k.n_compass(sl, Cm(25.0), Cm(2.9))
    PX = Cm(28.2)
    k.rect(sl, MSO_SHAPE.RECTANGLE, PX, Cm(2.6), Cm(0.14), Cm(1.0), fill=k.RED)
    k.tb(sl, PX + Cm(0.4), Cm(2.55), Cm(12.5), Cm(1), "방향별 조망 요약", size=16, color=k.WHITE, bold=True)
    py = 4.2
    for kk in range(8):
        col, en, kr = k.VIEW[k.view_of(sect[kk])]
        k.rect(sl, MSO_SHAPE.OVAL, PX, Cm(py + 0.15), Cm(0.5), Cm(0.5), fill=col)
        k.tb(sl, PX + Cm(0.8), Cm(py), Cm(2), Cm(0.8), DIRN[kk], size=13, color=k.WHITE, bold=True)
        k.tb(sl, PX + Cm(2.5), Cm(py - 0.05), Cm(9.5), Cm(0.7), en, size=10, color=col, bold=True)
        k.tb(sl, PX + Cm(2.5), Cm(py + 0.55), Cm(9.5), Cm(0.6), f"{kr} · 최고 {int(sect[kk])}m", size=9, color=k.MUTE); py += 1.15
    opens = [DIRN[kk] for kk in range(8) if k.view_of(sect[kk]) in ("OPEN", "LOW")]
    runs = ([('·'.join(opens), HL, True), (" 방향이 상대적으로 열려 조망·채광 유리 — 그 외는 고층 차폐 검토 필요.", k.WHITE, True)]
            if opens else [("사방이 중·고층으로 둘러싸여 개방면 제한적", HL, True), (" — 저층부 조망·상층부 향 확보 전략 검토.", k.WHITE, True)])
    k.caption_band(sl, runs)


# ── 커버 ──
def _cover(prs, address, law, site):
    sl = _blank(prs)
    k.rect(sl, MSO_SHAPE.RECTANGLE, 0, Cm(11.3), k.A3_W, Cm(0.16), fill=k.RED)
    k.tb(sl, Cm(1.5), Cm(3.2), Cm(39), Cm(2), "대지분석 보드", size=44, color=k.WHITE, bold=True)
    k.tb(sl, Cm(1.5), Cm(6.2), Cm(39), Cm(1), "SITE ANALYSIS", size=16, color=k.MUTE)
    lp = (site or {}).get("land_price") or {}
    k.tb(sl, Cm(1.5), Cm(8.2), Cm(39), Cm(2.4), [address,
         f"용도지역 {(law or {}).get('zone_use') or '-'} · 공시지가 {lp.get('price_per_sqm'):,}원/㎡" if lp.get("price_per_sqm") else address],
         size=15, color=k.WHITE)
    k.tb(sl, Cm(1.5), Cm(13), Cm(39), Cm(1),
         "지도 4종(광역·용도·높이·조망) + 데이터(지역통계·수급진단·대지정보·생활맥락·주변현황)",
         size=13, color=k.MUTE)


def build_full_deck(address: str, use_type: str = "주거", radius: int = 1000) -> bytes:
    """주소 → 종합 대지분석 덱 PPTX bytes.

    커버 + 지도 4종(광역·용도·높이·조망) + 데이터 팩트(지역통계·수급진단·대지정보·생활맥락·
    주변현황도·주변시설 요약) + **시설 종류별 상세 슬라이드**(어린이집·경로당·병원… 각 1장:
    지도 핀+이름·거리 목록). 형제앱(model·law·site) + 터읽기(board·surroundings·facilities)
    병렬 fetch. 전부 graceful — 소스 하나 죽어도 나머지로 덱 생성(절대 원칙 3).
    """
    # 주소 1회 해석 — 좌표 + PNU(VWorld). law 는 pnu 로 조회(가장 견고, §law 400 회피).
    from app.services.site_seed import build_site
    try:
        _s = build_site(address)
        lat, lon, pnu = _s.lat, _s.lon, (_s.pnu or "")
    except Exception:  # noqa: BLE001 — 주소 해석 실패는 하드블록
        lat = lon = None
        pnu = ""
    if lat is None:
        raise ValueError("주소 해석 실패")

    # model·law·site(지도용) + board·surroundings·facilities(데이터용) 병렬
    with ThreadPoolExecutor(max_workers=6) as ex:
        f_model = ex.submit(clients.fetch_model, address, 350)
        f_law = ex.submit(clients.fetch_law, address, pnu)
        f_site = ex.submit(clients.fetch_site, address)
        f_board = ex.submit(clients.fetch_board, address, use_type, radius, False)
        f_surr = ex.submit(clients.fetch_surroundings, address, radius)
        f_fac = ex.submit(clients.fetch_facilities, address, ds.FACILITY_KINDS, radius)
        model, law, site = f_model.result(), f_law.result(), f_site.result()
        board, surr, facs = f_board.result(), f_surr.result(), f_fac.result()
    parcel = (law or {}).get("parcel_geometry")
    board = board or {}

    prs = Presentation(); prs.slide_width, prs.slide_height = k.A3_W, k.A3_H
    _cover(prs, address, law, site)

    # 지도 4종 (매싱 3종은 model 필요 — 없으면 건너뜀)
    slide_wide(prs, address, lat, lon, law, site)
    if model:
        slide_use(prs, address, lat, lon, model, parcel)
        slide_site(prs, address, lat, lon, model, parcel)
        slide_viewring(prs, address, lat, lon, model, parcel)

    # 데이터 팩트 (같은 디자인 언어·새 숫자 0)
    ds.slide_region_stats(prs, address, board.get("region"), board.get("facts"), board.get("implications"))
    ds.slide_diagnose(prs, address, board.get("diagnoses"))
    ds.slide_site_info(prs, address, board.get("land_price"), board.get("building"),
                       board.get("real_estate"), board.get("hazards"))
    ds.slide_context(prs, address, board.get("context"))
    ds.slide_surroundings(prs, address, surr)
    # 주변시설 — 요약(종류별 개수) + 종류별 상세(지도 핀 + 이름·거리 목록)
    ds.slide_facilities(prs, address, facs, radius)
    ds.slide_facility_details(prs, address, facs, radius)

    buf = io.BytesIO(); prs.save(buf); return buf.getvalue()

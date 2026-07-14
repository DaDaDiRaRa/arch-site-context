"""C7 주변현황도 PPTX (CLAUDE.md §8.13, 심의 슬라이드 5).

SurroundingsResult → A3 1슬라이드: 반경 현황도(위성 + 반경밴드 + 카테고리별 색 점 + 범례)
+ 주변시설 카테고리 표 + 주변현황 서술문. 전부 네이티브(편집가능)·새 숫자 0.
"""

from __future__ import annotations

import io
from typing import Optional

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Pt

from app.schemas.surroundings import SurroundingsResult
from app.services import tiles

_A3_W, _A3_H = Cm(42.0), Cm(29.7)
_NAVY = RGBColor(0x2F, 0x54, 0x96)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_F = "맑은 고딕"
_MAP_PX = 1000


def _rgb(c):
    return RGBColor(c[0], c[1], c[2])


def build_surroundings_pptx(r: SurroundingsResult, client: Optional[httpx.Client] = None) -> bytes:
    own = client is None
    client = client or httpx.Client(timeout=20.0)
    try:
        prs = Presentation()
        prs.slide_width, prs.slide_height = _A3_W, _A3_H
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        tf = slide.shapes.add_textbox(Cm(1.2), Cm(0.5), Cm(39), Cm(1.3)).text_frame
        tf.text = "주변현황도"
        run = tf.paragraphs[0].runs[0]
        run.font.size, run.font.bold, run.font.name, run.font.color.rgb = Pt(22), True, _F, _NAVY

        # ── 왼쪽: 반경 현황도 ──
        BL, BT, BW = Cm(1.2), Cm(2.4), Cm(20.0)
        try:
            z = tiles.zoom_for_radius(r.site_lat, max(r.ring_radii or [r.radius]), _MAP_PX / 2 - 90)
            img, (cx, cy) = tiles.compose_basemap(r.site_lat, r.site_lon, z, _MAP_PX, _MAP_PX,
                                                  "vworld", client=client)
            buf = io.BytesIO()
            img.convert("RGB").save(buf, format="PNG")
            buf.seek(0)
            slide.shapes.add_picture(buf, BL, BT, BW, BW)
            scale = BW / _MAP_PX
            ox, oy = cx - _MAP_PX / 2, cy - _MAP_PX / 2
            ccx = ccy = _MAP_PX / 2

            def px(v):
                return Emu(int(v * scale))

            for rad in sorted(r.ring_radii or [r.radius]):
                rp = tiles.meters_to_pixels(rad, r.site_lat, z)
                ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(BL + (ccx - rp) * scale)),
                                            Emu(int(BT + (ccy - rp) * scale)), px(2 * rp), px(2 * rp))
                ov.fill.background()
                ov.line.color.rgb, ov.line.width = RGBColor(0xFF, 0xE0, 0x00), Pt(1.75)
                ov.name = f"반경 {rad}m"
            DOT = Cm(0.32)
            for cat in r.categories:
                for it in cat.items:
                    gx, gy = tiles.latlon_to_global_px(it.lat, it.lon, z)
                    x, y = gx - ox, gy - oy
                    if not (0 <= x <= _MAP_PX and 0 <= y <= _MAP_PX):
                        continue
                    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(BL + x * scale - DOT / 2)),
                                               Emu(int(BT + y * scale - DOT / 2)), DOT, DOT)
                    d.fill.solid()
                    d.fill.fore_color.rgb = _rgb(cat.color)
                    d.line.color.rgb, d.line.width = _WHITE, Pt(0.75)
                    d.name = f"{cat.name}: {it.name}"
            cm = Cm(0.5)
            c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(BL + ccx * scale - cm / 2)),
                                       Emu(int(BT + ccy * scale - cm / 2)), cm, cm)
            c.fill.solid()
            c.fill.fore_color.rgb, c.line.color.rgb, c.name = RGBColor(0xE5, 0, 0), _WHITE, "대상지"
        except Exception:
            slide.shapes.add_textbox(BL, BT, Cm(20), Cm(1)).text_frame.text = "(위성 지도 생성 실패)"

        # ── 오른쪽: 주변시설 카테고리 표 (범례 겸용 — 카테고리 셀에 색) ──
        cats = r.categories
        tbl = slide.shapes.add_table(len(cats) + 1, 3, Cm(22.0), Cm(2.4), Cm(18.5),
                                     Cm(1.0 * (len(cats) + 1))).table
        tbl.columns[0].width, tbl.columns[1].width, tbl.columns[2].width = Cm(4.0), Cm(2.5), Cm(12.0)
        for j, h in enumerate(["구분", "개수", "주요 시설"]):
            cell = tbl.cell(0, j)
            cell.text = h
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.size, p.runs[0].font.bold, p.runs[0].font.name = Pt(11), True, _F
            p.runs[0].font.color.rgb = _WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = _NAVY
        for i, cat in enumerate(cats, start=1):
            c0 = tbl.cell(i, 0)
            c0.text = cat.name
            c0.fill.solid()
            c0.fill.fore_color.rgb = _rgb(cat.color)
            p0 = c0.text_frame.paragraphs[0]
            p0.alignment = PP_ALIGN.CENTER
            p0.runs[0].font.size, p0.runs[0].font.bold, p0.runs[0].font.name = Pt(11), True, _F
            p0.runs[0].font.color.rgb = _WHITE
            c1 = tbl.cell(i, 1)
            c1.text = f"{cat.count}"
            c1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if c1.text_frame.paragraphs[0].runs:
                c1.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                c1.text_frame.paragraphs[0].runs[0].font.name = _F
            c2 = tbl.cell(i, 2)
            c2.text = ", ".join(it.name for it in cat.items[:5]) + (" 등" if cat.count > 5 else "")
            if c2.text_frame.paragraphs[0].runs:
                c2.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
                c2.text_frame.paragraphs[0].runs[0].font.name = _F

        # ── 하단: 주변현황 서술문 ──
        nb = slide.shapes.add_textbox(Cm(22.0), Cm(2.4) + Cm(1.0 * (len(cats) + 1)) + Cm(0.5),
                                      Cm(18.5), Cm(10.0)).text_frame
        nb.word_wrap = True
        nb.text = "■ 주변현황"
        nb.paragraphs[0].runs[0].font.size, nb.paragraphs[0].runs[0].font.bold = Pt(12), True
        nb.paragraphs[0].runs[0].font.name = _F
        p = nb.add_paragraph()
        p.text = r.narrative
        if p.runs:
            p.runs[0].font.size, p.runs[0].font.name = Pt(11), _F

        cap = slide.shapes.add_textbox(Cm(1.2), Cm(22.8), Cm(20), Cm(1.5)).text_frame
        cap.text = ("배경: VWorld 항공영상 · 시설: 카카오 검색(카테고리 코드 정제) · "
                    "반경밴드 " + "·".join(f"{x}m" for x in r.ring_radii) +
                    " · 도로폭·재개발경계는 소스 미확보(표기 안 함).")
        if cap.paragraphs[0].runs:
            cap.paragraphs[0].runs[0].font.size, cap.paragraphs[0].runs[0].font.name = Pt(9), _F

        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()
    finally:
        if own:
            client.close()

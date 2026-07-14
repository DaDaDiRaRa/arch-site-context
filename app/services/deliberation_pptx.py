"""심의 현황팩 PPTX 산출 (CLAUDE.md §8.13 C4·C5).

QuotaAssessment → A3 편집가능 PPT: 걸침 인구세대표 · 시설 현황표 · 편집가능 위치도
(위성 배경 + 네이티브 반경원·번호핀) · 총량제 판정박스. scratchpad 프로토타입의 앱 이식.

- 새 데이터·숫자 안 만듦 — assessment 를 렌더만 (절대 원칙 1·2).
- 표·도형은 네이티브(편집가능) — 사용자가 심의도서에 드롭 후 손질 ([[feedback-style]]).
- 위성 타일 실패해도 표는 나온다 (graceful). 폰트는 PowerPoint 가 렌더(파일엔 이름만).
"""

from __future__ import annotations

import io
from typing import List, Optional

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Pt

from app.schemas.quota import QuotaAssessment, QuotaResult
from app.schemas.survey import FacilityCategory
from app.services import tiles

_A3_W, _A3_H = Cm(42.0), Cm(29.7)
_NAVY = RGBColor(0x2F, 0x54, 0x96)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GREEN = RGBColor(0x2E, 0x7D, 0x32)
_RED = RGBColor(0xC6, 0x28, 0x28)
_HEAD = RGBColor(0xDD, 0xE3, 0xF0)
_F = "맑은 고딕"
_MAP_PX = 1000
_PIN_RGB = {"작은도서관": (0x21, 0x96, 0xF3), "경로당": (0xE9, 0x1E, 0x63),
            "어린이집": (0x4C, 0xAF, 0x50)}


def _blank(prs) -> "object":
    return prs.slides.add_slide(prs.slide_layouts[6])


def _title(slide, text: str) -> None:
    tf = slide.shapes.add_textbox(Cm(1.2), Cm(0.5), Cm(39), Cm(1.3)).text_frame
    tf.text = text
    r = tf.paragraphs[0].runs[0]
    r.font.size, r.font.bold, r.font.name, r.font.color.rgb = Pt(22), True, _F, _NAVY


def _setcell(cell, text, *, bold=False, fill=None, align=PP_ALIGN.CENTER,
             color=None, size=10) -> None:
    cell.text = "" if text is None else str(text)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    if p.runs:
        run = p.runs[0]
        run.font.size, run.font.bold, run.font.name = Pt(size), bold, _F
        if color:
            run.font.color.rgb = color
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def _survey_slide(prs, a: QuotaAssessment) -> None:
    slide = _blank(prs)
    _title(slide, "조사대상 행정동 인구·세대수 통계 (조사범위 걸침 합산)")
    hdr = ["구분", "총인구", "총세대", "비고(걸침율)", "적용인구", "적용세대"]
    rows = a.survey.dongs
    n = 1 + len(rows) + 1
    tbl = slide.shapes.add_table(n, 6, Cm(2.0), Cm(2.6), Cm(30), Cm(min(1.0, 22.0 / max(n, 1)) * n)).table
    for i, w in enumerate([Cm(6), Cm(5), Cm(5), Cm(5.5), Cm(4.5), Cm(4.5)]):
        tbl.columns[i].width = w
    for j, h in enumerate(hdr):
        _setcell(tbl.cell(0, j), h, bold=True, fill=_NAVY, color=_WHITE, size=12)
    for i, d in enumerate(rows, start=1):
        flag = " ⚠" if d.flagged else ""
        vals = [d.name + flag,
                f"{d.total_pop:,}" if d.total_pop else "확인필요",
                f"{d.total_hh:,}" if d.total_hh else "확인필요",
                f"{d.ratio*100:.2f}%",
                f"{d.applied_pop:,}" if d.applied_pop is not None else "-",
                f"{d.applied_hh:,}" if d.applied_hh is not None else "-"]
        for j, v in enumerate(vals):
            _setcell(tbl.cell(i, j), v, align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, size=11)
    tot = ["계", f"{a.survey.applied_pop_total:,}", f"{a.survey.applied_hh_total:,}",
           "", f"{a.survey.applied_pop_total:,}", f"{a.survey.applied_hh_total:,}"]
    for j, v in enumerate(tot):
        _setcell(tbl.cell(n - 1, j), v, bold=True, fill=_HEAD, size=11)
    ym_label = f"{a.ym[:4]}.{a.ym[4:]} 기준" if len(a.ym) == 6 else "기준월 미상"
    cap = slide.shapes.add_textbox(Cm(2.0), Cm(2.6) + Cm(1.0 * n) + Cm(0.2), Cm(30), Cm(1.2)).text_frame
    cap.text = (f"반경 {a.radius}m 걸침 행정동 · {ym_label}(행안부 주민등록) · "
                "걸침율=행정동 면적 대비 조사범위 교차비율 · ⚠타시군구 동은 '계' 제외(생활권 검토·사람 확정)")
    if cap.paragraphs[0].runs:
        cap.paragraphs[0].runs[0].font.size, cap.paragraphs[0].runs[0].font.name = Pt(10), _F


def _basemap(lat: float, lon: float, radius: int, client: httpx.Client):
    """VWorld 위성 배경 + 중심 전역픽셀. 실패 시 None."""
    try:
        z = tiles.zoom_for_radius(lat, radius, _MAP_PX / 2 - 100)
        img, (cx, cy) = tiles.compose_basemap(lat, lon, z, _MAP_PX, _MAP_PX, "vworld", client=client)
        return img.convert("RGB"), z, cx, cy
    except Exception:
        return None


def _facility_slide(prs, a: QuotaAssessment, cat: FacilityCategory,
                    client: httpx.Client) -> None:
    slide = _blank(prs)
    _title(slide, f"조사범위 내 {cat.category} 현황")
    pin_rgb = _PIN_RGB.get(cat.category, (0x60, 0x60, 0x60))

    # ── 왼쪽: 편집가능 위치도 (위성 배경 + 네이티브 반경원·번호핀) ──
    bm = _basemap(a.site_lat, a.site_lon, a.radius, client)
    if bm and a.site_lat:
        img, z, cx, cy = bm
        BL, BT, BW = Cm(1.2), Cm(2.4), Cm(19.5)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, BL, BT, BW, BW)
        scale = BW / _MAP_PX
        ox, oy = cx - _MAP_PX / 2, cy - _MAP_PX / 2

        def px(v):
            return Emu(int(v * scale))

        def to_xy(la, lo):
            gx, gy = tiles.latlon_to_global_px(la, lo, z)
            return gx - ox, gy - oy

        ccx = ccy = _MAP_PX / 2
        for r in sorted({500, a.radius}):
            rp = tiles.meters_to_pixels(r, a.site_lat, z)
            ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(BL + (ccx - rp) * scale)),
                                        Emu(int(BT + (ccy - rp) * scale)), px(2 * rp), px(2 * rp))
            ov.fill.background()
            ov.line.color.rgb, ov.line.width = RGBColor(0xFF, 0xE0, 0x00), Pt(2)
            ov.name = f"반경 {r}m"
        PR = Cm(0.62)
        for i, it in enumerate(cat.items, start=1):
            x, y = to_xy(it.lat, it.lon)
            if not (-30 <= x <= _MAP_PX + 30 and -30 <= y <= _MAP_PX + 30):
                continue
            pin = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(BL + x * scale - PR / 2)),
                                         Emu(int(BT + y * scale - PR / 2)), PR, PR)
            pin.fill.solid()
            pin.fill.fore_color.rgb = RGBColor(*pin_rgb)
            pin.line.color.rgb, pin.line.width = _WHITE, Pt(1.25)
            pin.name = f"{i}. {it.name}"
            tf = pin.text_frame
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(i)
            run.font.size, run.font.bold, run.font.name, run.font.color.rgb = Pt(10), True, _F, _WHITE
        cm = Cm(0.5)
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(BL + ccx * scale - cm / 2)),
                                   Emu(int(BT + ccy * scale - cm / 2)), cm, cm)
        c.fill.solid()
        c.fill.fore_color.rgb, c.line.color.rgb, c.name = _RED, _WHITE, "대상지"
        tx = Cm(21.5)
    else:
        slide.shapes.add_textbox(Cm(1.2), Cm(2.4), Cm(19), Cm(1)).text_frame.text = \
            "(위성 지도 생성 실패 — 표만)"
        tx = Cm(1.2)

    # ── 오른쪽(또는 아래): 현황 표 (번호=핀 번호) ──
    n = 1 + len(cat.items) + 1
    tbl = slide.shapes.add_table(n, 4, tx, Cm(2.4), Cm(19.5), Cm(min(0.62, 24.0 / max(n, 1)) * n)).table
    for i, w in enumerate([Cm(1.5), Cm(7.0), Cm(8.5), Cm(2.5)]):
        tbl.columns[i].width = w
    for j, h in enumerate(["번호", "이름", "주소", "면적"]):
        _setcell(tbl.cell(0, j), h, bold=True, fill=_NAVY, color=_WHITE, size=9)
    for i, it in enumerate(cat.items, start=1):
        _setcell(tbl.cell(i, 0), i, size=9)
        _setcell(tbl.cell(i, 1), it.name, align=PP_ALIGN.LEFT, size=9)
        _setcell(tbl.cell(i, 2), it.addr or "확인필요", align=PP_ALIGN.LEFT, size=9)
        _setcell(tbl.cell(i, 3), "확인필요", size=9)
    cap_txt = f"{cat.count} 개소"
    if cat.capacity:
        cap_txt += f" · 정원 {cat.capacity}명({cat.capacity_scope})"
    _setcell(tbl.cell(n - 1, 0), "계", bold=True, fill=_HEAD, size=9)
    _setcell(tbl.cell(n - 1, 1), cap_txt, bold=True, fill=_HEAD, align=PP_ALIGN.LEFT, size=9)
    _setcell(tbl.cell(n - 1, 2), "", fill=_HEAD)
    _setcell(tbl.cell(n - 1, 3), "", fill=_HEAD)


def _quota_slide(prs, res: QuotaResult, ctx: QuotaAssessment) -> None:
    slide = _blank(prs)
    label = f" [{res.label}]" if res.label else ""
    _title(slide, f"커뮤니티 설치계획 (총량제) 검토 — 부족/충족 판정{label}")
    sub = slide.shapes.add_textbox(Cm(1.2), Cm(1.9), Cm(39), Cm(0.8)).text_frame
    inf = f"{ctx.gu_infant:,}" if ctx.gu_infant else "확인필요"
    gu = f"{ctx.gu_households:,}" if ctx.gu_households else "확인필요"
    sub.text = (f"신축 {res.new_households:,}세대(설계) + 조사범위 적용 {res.applied_households:,}"
                f"세대(걸침) | 구 영유아 {inf} / 구 세대 {gu}")
    if sub.paragraphs[0].runs:
        sub.paragraphs[0].runs[0].font.size, sub.paragraphs[0].runs[0].font.name = Pt(10), _F

    xs, ys = [Cm(1.2), Cm(21.4)], [Cm(3.0), Cm(15.4)]
    for idx, f in enumerate(res.facilities):
        x, y = xs[idx % 2], ys[idx // 2]
        rowlabels = [("조례/법정 기준", f"법정 {f.legal_min:.0f}㎡ 이상" if f.legal_min else "필수 아님/미확정"),
                     ("산정 (자동)", _calc_text(f)),
                     ("판정 (자동)", f.verdict or "-"),
                     ("계획면적 (설계입력)", _plan_text(f))]
        tbl = slide.shapes.add_table(len(rowlabels) + 1, 2, x, y, Cm(19.2), Cm(11.0)).table
        tbl.columns[0].width, tbl.columns[1].width = Cm(4.5), Cm(14.7)
        hc = tbl.cell(0, 0)
        hc.merge(tbl.cell(0, 1))
        _setcell(hc, f"■ {f.name}", bold=True, fill=_NAVY, color=_WHITE,
                 align=PP_ALIGN.LEFT, size=13)
        for i, (lab, val) in enumerate(rowlabels, start=1):
            _setcell(tbl.cell(i, 0), lab, bold=True, fill=RGBColor(0xEE, 0xF1, 0xF8), size=10)
            col = None
            if lab.startswith("판정"):
                col = _RED if f.verdict == "부족시설" else (_GREEN if f.verdict == "충족시설" else None)
            _setcell(tbl.cell(i, 1), val, align=PP_ALIGN.LEFT, color=col,
                     bold=lab.startswith("판정"), size=10)
    cap = slide.shapes.add_textbox(Cm(1.2), Cm(27.6), Cm(39), Cm(1.5)).text_frame
    cap.text = ("공식: 서울시 주민공동시설 총량제·조례(community_quota.json) · 산출>0=부족·≤0=충족 · "
                "법정면적은 조례 변동값(confidence 낮으면 조례 확인) · 판정은 '참고', 최종 확정은 사람.")
    if cap.paragraphs[0].runs:
        cap.paragraphs[0].runs[0].font.size, cap.paragraphs[0].runs[0].font.name = Pt(9), _F


def _calc_text(f) -> str:
    parts: List[str] = []
    if f.expected_people is not None:
        parts.append(f"예상인원 = {f.expected_people:,.2f}명")
    if f.required_area is not None:
        parts.append(f"산출면적 = {f.required_area:,.2f}㎡ (기존 {f.existing_area:,.1f}㎡ 차감)")
    for nt in f.notes:
        if "조례" in nt:
            parts.append(f"※ {nt}")
    return "\n".join(parts) if parts else "면적 기준(고정)"


def _plan_text(f) -> str:
    if f.planned_area is None:
        return "계획면적 미입력"
    if f.legal_min is None:
        return f"계획 {f.planned_area:,.2f}㎡ (법정 미확정)"
    ok = "충족 ✓" if f.plan_ok else "미달 ✗"
    return f"계획 {f.planned_area:,.2f}㎡ vs 법정 {f.legal_min:.0f}㎡ → {ok} ({f.plan_diff:+,.2f}㎡)"


def build_pptx(a: QuotaAssessment, client: Optional[httpx.Client] = None) -> bytes:
    """QuotaAssessment → A3 편집가능 심의 현황팩 PPTX (bytes)."""
    own = client is None
    client = client or httpx.Client(timeout=20.0)
    try:
        prs = Presentation()
        prs.slide_width, prs.slide_height = _A3_W, _A3_H
        _survey_slide(prs, a)
        for cat in a.facilities:
            if cat.category in ("작은도서관", "경로당") and cat.items:
                _facility_slide(prs, a, cat, client)
        for res in a.results:
            _quota_slide(prs, res, a)
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    finally:
        if own:
            client.close()
